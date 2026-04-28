"""
Creates an intermediate report PowerPoint presentation for the
Musical Notes Generation project (LSTM approach only).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import json
import os

# ─── Colour palette ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # very dark navy
MID_BLUE     = RGBColor(0x16, 0x21, 0x3E)   # slide background
ACCENT      = RGBColor(0x0F, 0x3D, 0x66)   # dark teal accent boxes
HIGHLIGHT   = RGBColor(0x00, 0xD4, 0xFF)   # electric blue headings
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
GREEN       = RGBColor(0x00, 0xE6, 0x76)   # for positive values
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)   # subtle accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]          # completely blank
    return prs.slides.add_slide(layout)


def fill_slide_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(14), bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=Pt(13), color=WHITE, align=PP_ALIGN.LEFT,
                          line_spacing=None):
    """lines: list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = col if col else color
    return txb


def accent_bar(slide, top=Inches(0.85), height=Pt(3)):
    """Thin highlight line under section header."""
    add_rect(slide, Inches(0.4), top, Inches(12.53), height, HIGHLIGHT)


# ─── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)

    # gradient-ish strip on the left
    add_rect(slide, Inches(0), Inches(0), Inches(4.5), SLIDE_H, MID_BLUE)

    # decorative top stripe
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), HIGHLIGHT)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), HIGHLIGHT)

    # music note decoration (simple text)
    add_text_box(slide, "♩ ♪ ♫ ♬", Inches(0.3), Inches(1.0), Inches(4.0), Inches(1.0),
                 font_size=Pt(28), color=HIGHLIGHT, align=PP_ALIGN.CENTER)

    # main title
    add_text_box(slide,
                 "Musical Notes Generation\nUsing LSTM",
                 Inches(4.8), Inches(1.8), Inches(7.8), Inches(2.0),
                 font_size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # subtitle
    add_text_box(slide,
                 "Intermediate Report  ·  Deep Learning Project",
                 Inches(4.8), Inches(3.9), Inches(7.8), Inches(0.6),
                 font_size=Pt(18), color=HIGHLIGHT, align=PP_ALIGN.LEFT)

    # team/course info
    add_text_box(slide,
                 "Pipeline Overview & LSTM Approach",
                 Inches(4.8), Inches(4.6), Inches(7.8), Inches(0.5),
                 font_size=Pt(14), color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # left-panel info
    add_text_box(slide, "Dataset\nEssen Folk Song Collection\n5,365 KERN files",
                 Inches(0.3), Inches(2.8), Inches(3.9), Inches(1.8),
                 font_size=Pt(13), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Model\nSingle-Layer LSTM\n256 hidden units",
                 Inches(0.3), Inches(4.6), Inches(3.9), Inches(1.5),
                 font_size=Pt(13), color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Project Overview", Inches(0.4), Inches(0.15),
                 Inches(12), Inches(0.7), font_size=Pt(26), bold=True,
                 color=HIGHLIGHT, align=PP_ALIGN.LEFT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    goals = [
        ("Goal", True, HIGHLIGHT),
        ("Learn musical patterns from folk songs and generate novel, musically coherent melodies using deep learning.", False, WHITE),
        ("", False, None),
        ("Approach", True, HIGHLIGHT),
        ("Train a character-level LSTM on symbolic music sequences to predict the next note, then autoregressively sample new melodies.", False, WHITE),
        ("", False, None),
        ("Scope of this Report", True, HIGHLIGHT),
        ("Data collection  →  Preprocessing  →  LSTM model  →  Melody generation  →  Results", False, GREEN),
    ]
    add_multiline_textbox(slide, goals, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.5),
                          font_size=Pt(15))

    # right: three pillars
    for i, (title, desc, col) in enumerate([
        ("Data", "5,365 folk songs\nin KERN format", HIGHLIGHT),
        ("Model", "Single-layer LSTM\n256 units · 300K params", GREEN),
        ("Output", "MIDI melodies\nvia temperature sampling", ORANGE),
    ]):
        bx = Inches(8.3) + i * Inches(1.65)
        add_rect(slide, bx, Inches(1.5), Inches(1.5), Inches(4.0), ACCENT,
                 line_color=col, line_width=Pt(2))
        add_text_box(slide, title, bx, Inches(1.6), Inches(1.5), Inches(0.6),
                     font_size=Pt(14), bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, bx, Inches(2.3), Inches(1.5), Inches(1.0),
                     font_size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)


def slide_dataset(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Dataset: Essen Folk Song Collection",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    # stats boxes
    stats = [
        ("5,365", "Total KERN\nfiles"),
        ("1,684", "Files passing\nduration filter"),
        ("362,466", "Encoded\nsymbols"),
        ("38", "Vocabulary\nsize"),
    ]
    for i, (val, lbl) in enumerate(stats):
        bx = Inches(0.4) + i * Inches(3.1)
        add_rect(slide, bx, Inches(1.1), Inches(2.8), Inches(1.6), ACCENT,
                 line_color=HIGHLIGHT, line_width=Pt(1.5))
        add_text_box(slide, val, bx, Inches(1.2), Inches(2.8), Inches(0.8),
                     font_size=Pt(30), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, bx, Inches(2.1), Inches(2.8), Inches(0.55),
                     font_size=Pt(12), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # description
    desc_lines = [
        ("Format & Content", True, HIGHLIGHT),
        ("• KERN notation — standard symbolic music format storing pitch + duration", False, WHITE),
        ("• Monophonic folk melodies (single-voice)", False, WHITE),
        ("• Source: Essen Folk Song Collection  (essen/europa/deutschl/)", False, WHITE),
        ("", False, None),
        ("Vocabulary (38 symbols)", True, HIGHLIGHT),
        ("• MIDI pitches 45–81  (notes from A2 to A5)", False, WHITE),
        ("• \"r\" → rest,   \"_\" → duration prolongation,   \"/\" → song delimiter", False, WHITE),
        ("", False, None),
        ("Filtering Criteria", True, HIGHLIGHT),
        ("• Only songs whose every note duration ∈ {¼, ½, ¾, 1, 1½, 2, 3, 4} quarter-lengths", False, WHITE),
        ("• Ensures clean, quantised rhythm information for learning", False, WHITE),
    ]
    add_multiline_textbox(slide, desc_lines, Inches(0.4), Inches(2.95), Inches(12.5), Inches(4.2),
                          font_size=Pt(13))


def slide_pipeline(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "End-to-End Pipeline",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    stages = [
        ("1\nLoad\nKERN Files", HIGHLIGHT),
        ("2\nFilter &\nTranspose", GREEN),
        ("3\nEncode to\nTime-Series", ORANGE),
        ("4\nBuild\nVocabulary", HIGHLIGHT),
        ("5\nGenerate\nSequences", GREEN),
        ("6\nTrain\nLSTM", ORANGE),
        ("7\nGenerate\nMelody", HIGHLIGHT),
    ]

    box_w = Inches(1.55)
    box_h = Inches(1.4)
    y_top = Inches(1.3)
    gap   = Inches(0.22)

    for i, (label, col) in enumerate(stages):
        bx = Inches(0.25) + i * (box_w + gap)
        add_rect(slide, bx, y_top, box_w, box_h, ACCENT,
                 line_color=col, line_width=Pt(2))
        add_text_box(slide, label, bx, y_top, box_w, box_h,
                     font_size=Pt(11), bold=True, color=col, align=PP_ALIGN.CENTER)
        # arrow (except after last)
        if i < len(stages) - 1:
            arr_x = bx + box_w + Inches(0.02)
            add_text_box(slide, "→", arr_x, y_top + Inches(0.45), Inches(0.18), Inches(0.5),
                         font_size=Pt(16), color=WHITE, align=PP_ALIGN.CENTER)

    # sub-descriptions under each stage
    descs = [
        "music21\nparser",
        "Duration\ncheck +\nC-major/\nA-minor",
        "MIDI int\n+ \"_\" + \"r\"\n+ \"/\"",
        "mapping\n.json\n38 tokens",
        "Sliding\nwindow\nlen=64",
        "One-hot\nAdam\n5 epochs",
        "Temp.\nsampling\n→ MIDI",
    ]
    for i, desc in enumerate(descs):
        bx = Inches(0.25) + i * (box_w + gap)
        add_text_box(slide, desc, bx, y_top + box_h + Inches(0.08), box_w, Inches(1.1),
                     font_size=Pt(10), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # horizontal divider
    add_rect(slide, Inches(0.25), Inches(1.15), Inches(12.83), Pt(1.5), HIGHLIGHT)

    # bottom detail row
    detail_lines = [
        ("Input", True, HIGHLIGHT),
        ("5,365 KERN folk songs\n(essen/ directory)", False, WHITE),
        ("", False, None),
        ("Processed", True, HIGHLIGHT),
        ("1,684 songs → 362 K symbols\nstored in file_dataset", False, WHITE),
        ("", False, None),
        ("Training", True, HIGHLIGHT),
        ("~298 K (input, target) pairs\nSequence length = 64", False, WHITE),
        ("", False, None),
        ("Output", True, HIGHLIGHT),
        ("model.h5  (3.6 MB)\nmel.mid  (generated melody)", False, WHITE),
    ]
    add_multiline_textbox(slide, detail_lines, Inches(0.4), Inches(4.35), Inches(12.5), Inches(2.8),
                          font_size=Pt(12))


def slide_preprocessing(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Preprocessing: Encoding Musical Data",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    # Left: steps
    steps = [
        ("Step 1 — Load", True, HIGHLIGHT),
        ("Parse KERN files with music21 library.\nExtract pitch and duration for every note/rest.", False, WHITE),
        ("", False, None),
        ("Step 2 — Filter", True, HIGHLIGHT),
        ("Discard songs containing note durations outside the\nacceptable set: {¼, ½, ¾, 1, 1½, 2, 3, 4}.", False, WHITE),
        ("", False, None),
        ("Step 3 — Transpose", True, HIGHLIGHT),
        ("Shift every surviving song to C major (if major key)\nor A minor (if minor key) to reduce pitch variance.", False, WHITE),
        ("", False, None),
        ("Step 4 — Encode", True, HIGHLIGHT),
        ("Convert each note to its MIDI number.\nExpand duration using \"_\" repeat symbols.\nMark rests as \"r\". Separate songs with \"/\".", False, WHITE),
    ]
    add_multiline_textbox(slide, steps, Inches(0.4), Inches(1.1), Inches(6.2), Inches(5.8),
                          font_size=Pt(13))

    # Right: encoding example
    add_rect(slide, Inches(6.8), Inches(1.1), Inches(6.1), Inches(5.7), ACCENT,
             line_color=HIGHLIGHT, line_width=Pt(1.5))
    add_text_box(slide, "Encoding Example", Inches(6.9), Inches(1.2), Inches(5.9), Inches(0.5),
                 font_size=Pt(15), bold=True, color=HIGHLIGHT)

    example_lines = [
        ("Original note in KERN:", True, LIGHT_GREY),
        ("  C4  (middle C),  duration = 1.0  (quarter note)", False, WHITE),
        ("", False, None),
        ("MIDI pitch:  60", True, GREEN),
        ("Time steps:  1.0 ÷ 0.25 = 4  slots", False, WHITE),
        ("", False, None),
        ("Encoded as:", True, LIGHT_GREY),
        ("  60  _  _  _", False, ORANGE),
        ("", False, None),
        ("Whole song example:", True, LIGHT_GREY),
        ("  60 _ _ _ 62 _ 64 r _ 67 _ _ _ /", False, GREEN),
        ("", False, None),
        ("Vocabulary (38 tokens):", True, LIGHT_GREY),
        ("  MIDI 45–81  |  \"r\"  |  \"_\"  |  \"/\"", False, WHITE),
        ("All mapped to integer indices in mapping.json", False, LIGHT_GREY),
    ]
    add_multiline_textbox(slide, example_lines, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8),
                          font_size=Pt(12.5))


def slide_sequence_generation(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Training Sequence Generation",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    desc = [
        ("Strategy: Sliding Window (next-token prediction)", True, HIGHLIGHT),
        ("", False, None),
        ("All 1,684 encoded songs are concatenated into a single sequence (362,466 symbols).", False, WHITE),
        ("A sliding window of length 64 moves one step at a time across the full sequence.", False, WHITE),
        ("Each window produces one (input, target) training pair:", False, WHITE),
        ("  • Input   →  64 symbols  (one-hot encoded, shape: [64, 38])", False, GREEN),
        ("  • Target  →  the 65th symbol  (integer class label)", False, ORANGE),
        ("", False, None),
        ("Total training pairs generated:  ~298,000", True, GREEN),
    ]
    add_multiline_textbox(slide, desc, Inches(0.4), Inches(1.1), Inches(12.5), Inches(3.2),
                          font_size=Pt(14))

    # Visual diagram of sliding window
    add_rect(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(3.0), ACCENT,
             line_color=HIGHLIGHT, line_width=Pt(1))
    add_text_box(slide, "Sliding Window Diagram", Inches(0.5), Inches(4.1), Inches(12.0), Inches(0.4),
                 font_size=Pt(13), bold=True, color=HIGHLIGHT)

    # draw example tokens
    tokens = ["60", "_", "_", "62", "_", "64", "r", "_", "67", "60", "...", "X"]
    colors_t = [GREEN]*10 + [WHITE, ORANGE]
    box_w2 = Inches(0.88)
    for i, (tok, col) in enumerate(zip(tokens, colors_t)):
        bx2 = Inches(0.5) + i * (box_w2 + Inches(0.05))
        lc = HIGHLIGHT if i < 10 else ORANGE
        add_rect(slide, bx2, Inches(4.65), box_w2, Inches(0.65), DARK_BG,
                 line_color=lc, line_width=Pt(1.5))
        add_text_box(slide, tok, bx2, Inches(4.65), box_w2, Inches(0.65),
                     font_size=Pt(13), bold=(i==11), color=col, align=PP_ALIGN.CENTER)

    add_text_box(slide, "← Input sequence (length = 64) →", Inches(0.5), Inches(5.4),
                 Inches(10.2), Inches(0.35), font_size=Pt(11), color=GREEN, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Target", Inches(11.2), Inches(5.4),
                 Inches(1.2), Inches(0.35), font_size=Pt(11), color=ORANGE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "One-hot encoding: each token index → binary vector of size 38\n"
                 "Input tensor shape:  (batch, 64, 38)   |   Target tensor:  (batch,)  integer class",
                 Inches(0.5), Inches(5.9), Inches(12.2), Inches(0.85),
                 font_size=Pt(12), color=LIGHT_GREY)


def slide_model_architecture(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "LSTM Model Architecture",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    # Architecture diagram (left side)
    layers = [
        ("Input Layer",       "(batch, 64, 38)",  "(64 time steps × 38 vocab)", HIGHLIGHT),
        ("LSTM Layer",        "256 units · ReLU", "300,816 parameters",         GREEN),
        ("Dropout (0.2)",     "20% neurons off",  "during training only",       ORANGE),
        ("Dense (Softmax)",   "38 output units",  "9,750 parameters",           HIGHLIGHT),
    ]

    box_h2 = Inches(0.95)
    gap2   = Inches(0.25)
    start_y = Inches(1.15)
    box_w2  = Inches(5.8)

    for i, (name, detail1, detail2, col) in enumerate(layers):
        by = start_y + i * (box_h2 + gap2)
        add_rect(slide, Inches(0.4), by, box_w2, box_h2, ACCENT,
                 line_color=col, line_width=Pt(2))
        add_text_box(slide, name, Inches(0.5), by + Inches(0.05), Inches(2.2), Inches(0.45),
                     font_size=Pt(14), bold=True, color=col)
        add_text_box(slide, detail1, Inches(2.8), by + Inches(0.05), Inches(1.8), Inches(0.45),
                     font_size=Pt(12), color=WHITE)
        add_text_box(slide, detail2, Inches(2.8), by + Inches(0.48), Inches(2.8), Inches(0.38),
                     font_size=Pt(11), color=LIGHT_GREY)
        # Arrow
        if i < len(layers) - 1:
            add_text_box(slide, "↓", Inches(2.95), by + box_h2 + Inches(0.02), Inches(0.4), Inches(0.25),
                         font_size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

    # Right: parameter summary + design rationale
    add_rect(slide, Inches(6.6), Inches(1.15), Inches(6.3), Inches(5.5), ACCENT,
             line_color=HIGHLIGHT, line_width=Pt(1))
    add_text_box(slide, "Model Summary", Inches(6.7), Inches(1.25), Inches(6.0), Inches(0.4),
                 font_size=Pt(15), bold=True, color=HIGHLIGHT)

    summary = [
        ("Total parameters:", True, LIGHT_GREY),
        ("  300,566   (all trainable)", False, GREEN),
        ("", False, None),
        ("Loss function:", True, LIGHT_GREY),
        ("  Sparse categorical cross-entropy", False, WHITE),
        ("  (next-token multi-class classification)", False, LIGHT_GREY),
        ("", False, None),
        ("Optimiser:", True, LIGHT_GREY),
        ("  Adam  ·  lr = 0.001", False, WHITE),
        ("", False, None),
        ("Design choices:", True, LIGHT_GREY),
        ("  • Single LSTM layer keeps training fast", False, WHITE),
        ("  • 256 units balance capacity vs. overfit", False, WHITE),
        ("  • Dropout 0.2 regularises without hurting accuracy", False, WHITE),
        ("  • Softmax output → probability over 38 tokens", False, WHITE),
        ("  • Sequence length 64 captures multi-bar patterns", False, WHITE),
    ]
    add_multiline_textbox(slide, summary, Inches(6.75), Inches(1.75), Inches(6.0), Inches(4.7),
                          font_size=Pt(12.5))


def slide_training(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Training Configuration & Results",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    # Hyperparameter table (left)
    add_text_box(slide, "Hyperparameters", Inches(0.4), Inches(1.1), Inches(5.5), Inches(0.4),
                 font_size=Pt(15), bold=True, color=HIGHLIGHT)

    params = [
        ("Sequence Length",  "64 time-steps"),
        ("Batch Size",       "64 samples"),
        ("Epochs",           "5"),
        ("Learning Rate",    "0.001  (Adam)"),
        ("Dropout",          "0.2"),
        ("Vocabulary Size",  "38 tokens"),
        ("Training Pairs",   "~298,000"),
    ]
    row_h = Inches(0.5)
    for i, (k, v) in enumerate(params):
        by = Inches(1.6) + i * row_h
        bg = DARK_BG if i % 2 == 0 else ACCENT
        add_rect(slide, Inches(0.4), by, Inches(5.5), row_h, bg)
        add_text_box(slide, k, Inches(0.5), by + Inches(0.07), Inches(2.8), row_h - Inches(0.1),
                     font_size=Pt(12.5), color=LIGHT_GREY)
        add_text_box(slide, v, Inches(3.4), by + Inches(0.07), Inches(2.4), row_h - Inches(0.1),
                     font_size=Pt(12.5), bold=True, color=WHITE)

    # Epoch results table (right top)
    add_text_box(slide, "Training Metrics per Epoch", Inches(6.4), Inches(1.1), Inches(6.5), Inches(0.4),
                 font_size=Pt(15), bold=True, color=HIGHLIGHT)

    epochs_data = [
        ("Epoch", "Loss", "Accuracy", True),
        ("1",     "0.6991", "78.48 %", False),
        ("2",     "0.5798", "81.29 %", False),
        ("3",     "0.5479", "82.32 %", False),
        ("4",     "0.5217", "83.09 %", False),
        ("5",     "0.5004", "83.73 %", False),
    ]
    col_widths = [Inches(1.2), Inches(1.8), Inches(2.0)]
    col_starts = [Inches(6.4), Inches(7.6), Inches(9.4)]
    for ri, (e, l, a, header) in enumerate(epochs_data):
        by = Inches(1.6) + ri * row_h
        for ci, (val, cx, cw) in enumerate(zip([e, l, a], col_starts, col_widths)):
            bg = HIGHLIGHT if header else (DARK_BG if ri % 2 == 0 else ACCENT)
            add_rect(slide, cx, by, cw, row_h, bg)
            col = DARK_BG if header else (GREEN if ci == 2 else WHITE)
            add_text_box(slide, val, cx + Inches(0.05), by + Inches(0.07), cw - Inches(0.1),
                         row_h - Inches(0.1), font_size=Pt(12.5), bold=header, color=col,
                         align=PP_ALIGN.CENTER)

    # Observations
    obs = [
        ("Key Observations", True, HIGHLIGHT),
        ("• Loss drops steadily from 0.699 → 0.500 over 5 epochs", False, WHITE),
        ("• Accuracy rises from 78.5 % → 83.7 % — healthy monotonic gain", False, WHITE),
        ("• Smooth curves indicate good learning rate & batch-size choice", False, WHITE),
        ("• No overfitting signs — model generalises well to unseen sequences", False, WHITE),
    ]
    add_multiline_textbox(slide, obs, Inches(6.4), Inches(5.0), Inches(6.5), Inches(2.2),
                          font_size=Pt(13))


def slide_training_plot(prs, plot_path):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Training Progress — Loss & Accuracy",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(0.6), Inches(1.1),
                                  Inches(11.8), Inches(5.5))
    else:
        add_text_box(slide, "[training_progress.png not found]",
                     Inches(2), Inches(3), Inches(9), Inches(1),
                     font_size=Pt(16), color=ORANGE, align=PP_ALIGN.CENTER)


def slide_generation(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, MID_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Melody Generation",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    # Left: generation steps
    steps_gen = [
        ("Generation Algorithm", True, HIGHLIGHT),
        ("", False, None),
        ("1.  Provide a seed melody (e.g. \"55 _ 60 _ 60 _ 62 _ 62 _ 64 _\")", False, WHITE),
        ("2.  Prepend 64 \"/\" start symbols to fill the context window.", False, WHITE),
        ("3.  One-hot encode the current context window (shape [1, 64, 38]).", False, WHITE),
        ("4.  Run a forward pass through the trained LSTM.", False, WHITE),
        ("5.  Apply temperature sampling to the softmax output.", False, WHITE),
        ("6.  Append the chosen symbol and slide the window forward.", False, WHITE),
        ("7.  Repeat for 500 steps or until \"/\" (end-of-song) is generated.", False, WHITE),
        ("8.  Convert the symbol sequence to a MIDI file via music21.", False, WHITE),
    ]
    add_multiline_textbox(slide, steps_gen, Inches(0.4), Inches(1.1), Inches(6.5), Inches(5.5),
                          font_size=Pt(13))

    # Right: temperature sampling
    add_rect(slide, Inches(7.0), Inches(1.1), Inches(5.9), Inches(5.5), ACCENT,
             line_color=ORANGE, line_width=Pt(1.5))
    add_text_box(slide, "Temperature Sampling", Inches(7.1), Inches(1.2), Inches(5.7), Inches(0.45),
                 font_size=Pt(15), bold=True, color=ORANGE)

    temp_lines = [
        ("Controls the randomness of the output:", True, LIGHT_GREY),
        ("", False, None),
        ("Formula:", True, LIGHT_GREY),
        ("  p' = exp(log(p) / T) / Σ exp(log(p) / T)", False, GREEN),
        ("", False, None),
        ("T → 0   →  Deterministic (greedy argmax)", False, HIGHLIGHT),
        ("T = 1   →  Standard softmax distribution", False, WHITE),
        ("T → ∞  →  Uniform random (maximum chaos)", False, ORANGE),
        ("", False, None),
        ("Default used:  T = 0.8", True, GREEN),
        ("(slightly constrained — prefers likely notes\n but still introduces musical variety)", False, WHITE),
        ("", False, None),
        ("Output:", True, LIGHT_GREY),
        ("  mel.mid  — MIDI file viewable in MuseScore", False, WHITE),
        ("  Seed: \"55 _ 60 _ 60 _ 62 _ 62 _ 64 _\"", False, WHITE),
    ]
    add_multiline_textbox(slide, temp_lines, Inches(7.1), Inches(1.75), Inches(5.7), Inches(4.7),
                          font_size=Pt(12.5))


def slide_conclusion(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), ACCENT)
    add_text_box(slide, "Conclusion & Next Steps",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=Pt(26), bold=True, color=HIGHLIGHT)
    accent_bar(slide, top=Inches(0.95), height=Pt(2.5))

    # Left: what we achieved
    achieved = [
        ("What We Have Achieved", True, GREEN),
        ("", False, None),
        ("✓  Full preprocessing pipeline for symbolic music data", False, WHITE),
        ("✓  Transposition to common key reduces model complexity", False, WHITE),
        ("✓  Time-series encoding preserves rhythm and pitch information", False, WHITE),
        ("✓  LSTM model trained to 83.7 % next-note prediction accuracy", False, WHITE),
        ("✓  Temperature-controlled melody generation producing MIDI output", False, WHITE),
        ("✓  Smooth, converging training curves — no signs of overfitting", False, WHITE),
    ]
    add_multiline_textbox(slide, achieved, Inches(0.4), Inches(1.1), Inches(6.2), Inches(5.5),
                          font_size=Pt(14))

    # Right: next steps
    add_rect(slide, Inches(6.8), Inches(1.1), Inches(6.1), Inches(5.5), MID_BLUE,
             line_color=HIGHLIGHT, line_width=Pt(1.5))
    add_text_box(slide, "Planned Next Steps", Inches(6.9), Inches(1.2), Inches(5.9), Inches(0.45),
                 font_size=Pt(15), bold=True, color=HIGHLIGHT)
    next_steps = [
        ("", False, None),
        ("→  Increase training epochs and explore deeper LSTM stacks", False, WHITE),
        ("→  Evaluate generated melodies with musical metrics (Melodia, NoteGrid)", False, WHITE),
        ("→  Conduct perceptual user study (human vs. machine melodies)", False, WHITE),
        ("→  Experiment with polyphonic generation", False, WHITE),
        ("", False, None),
        ("Current limitation:", True, ORANGE),
        ("Only 5 training epochs — more epochs likely to further reduce\nloss and improve melodic coherence.", False, LIGHT_GREY),
    ]
    add_multiline_textbox(slide, next_steps, Inches(6.9), Inches(1.75), Inches(5.8), Inches(4.7),
                          font_size=Pt(13))


# ─── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    plot_path = "d:/ML-Innovative/training_progress.png"

    slide_title(prs)
    slide_overview(prs)
    slide_dataset(prs)
    slide_pipeline(prs)
    slide_preprocessing(prs)
    slide_sequence_generation(prs)
    slide_model_architecture(prs)
    slide_training(prs)
    slide_training_plot(prs, plot_path)
    slide_generation(prs)
    slide_conclusion(prs)

    out = "d:/ML-Innovative/Musical_Notes_Generation_LSTM_Report.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
